from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ============================================================
# カラーパレット
# ============================================================
C_DARK   = RGBColor(0x14, 0x0A, 0x04)   # #140A04 ダークブラウン（ヘッダー）
C_ACCENT = RGBColor(0xE8, 0x7A, 0x2A)   # #E87A2A オレンジ（アクセント）
C_LIGHT  = RGBColor(0xFF, 0xF8, 0xF0)   # #FFF8F0 クリーム（背景）
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY   = RGBColor(0x66, 0x60, 0x5A)
C_SUBTLE = RGBColor(0xF5, 0xED, 0xE0)   # 薄いベージュ

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # 完全ブランク

# ============================================================
# ヘルパー関数
# ============================================================
def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape

def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=C_DARK,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Hiragino Sans"
    return txBox

def add_multiline(slide, lines, x, y, w, h,
                  font_size=16, color=C_DARK, line_spacing=1.15):
    """lines: list of (text, bold, font_size_override or None)"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, fs = item, False, font_size
        elif len(item) == 2:
            text, bold = item; fs = font_size
        else:
            text, bold, fs = item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fs)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Hiragino Sans"
    return txBox

def bg_cream(slide):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_LIGHT)

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.3), C_DARK)
    add_textbox(slide, title,
                Inches(0.5), Inches(0.18), Inches(11), Inches(0.7),
                font_size=30, bold=True, color=C_WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.5), Inches(0.82), Inches(11), Inches(0.4),
                    font_size=14, color=C_ACCENT)

def accent_line(slide, y=Inches(1.3)):
    add_rect(slide, 0, y, SLIDE_W, Inches(0.06), C_ACCENT)

def card(slide, x, y, w, h, fill=C_WHITE):
    shape = add_rect(slide, x, y, w, h, fill)
    shape.line.color.rgb = RGBColor(0xDD, 0xCC, 0xBB)
    shape.line.width = Pt(0.75)
    return shape

# ============================================================
# スライド 1: タイトル
# ============================================================
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, C_DARK)
# オレンジの帯（下部）
add_rect(sl, 0, Inches(5.8), SLIDE_W, Inches(1.7), C_ACCENT)

add_textbox(sl, "3rd Place Hub",
            Inches(1), Inches(1.2), Inches(11), Inches(1.4),
            font_size=60, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_textbox(sl, "メルボルン日本人コミュニティ × 仕事・住居マッチング",
            Inches(1), Inches(2.7), Inches(11), Inches(0.7),
            font_size=22, color=C_ACCENT, align=PP_ALIGN.CENTER)

add_textbox(sl, "先輩の経験を、後輩へ。",
            Inches(1), Inches(3.5), Inches(11), Inches(0.6),
            font_size=20, color=C_LIGHT, italic=True, align=PP_ALIGN.CENTER)

add_textbox(sl, "運営メンバー説明会　2026",
            Inches(1), Inches(6.0), Inches(11), Inches(0.5),
            font_size=16, color=C_DARK, align=PP_ALIGN.CENTER)

# ============================================================
# スライド 2: アジェンダ
# ============================================================
sl = prs.slides.add_slide(blank_layout)
bg_cream(sl)
header_bar(sl, "本日のアジェンダ")
accent_line(sl)

items = [
    "01  3rd Place Hub とは何か",
    "02  なぜ今メルボルンで必要なのか",
    "03  サービスの仕組み",
    "04  コミュニティの現在地",
    "05  ビジネスモデル",
    "06  これからのビジョンとフェーズ計画",
    "07  一緒に運営しませんか",
]
for i, item in enumerate(items):
    y = Inches(1.55) + i * Inches(0.74)
    card(sl, Inches(1.5), y, Inches(10), Inches(0.62), C_WHITE)
    add_textbox(sl, item, Inches(1.8), y + Pt(5), Inches(9.5), Inches(0.55),
                font_size=18, color=C_DARK)

# ============================================================
# スライド 3: 3rd Place Hub とは
# ============================================================
sl = prs.slides.add_slide(blank_layout)
bg_cream(sl)
header_bar(sl, "3rd Place Hub とは", "メルボルン日本人コミュニティ発・仕事と住居のマッチングプラットフォーム")
accent_line(sl)

# 大見出しキャッチ
add_rect(sl, Inches(0.5), Inches(1.55), Inches(12.33), Inches(1.1), C_ACCENT)
add_textbox(sl,
            "先輩の「生の経験」を、これから渡航する後輩へバトンパスする",
            Inches(0.7), Inches(1.62), Inches(12), Inches(0.9),
            font_size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 3カラム
col_x = [Inches(0.4), Inches(4.65), Inches(8.9)]
col_titles = ["第1の場所\n（家）", "第2の場所\n（職場・学校）", "第3の場所\n（コミュニティ）"]
col_texts  = [
    "自分の部屋・住まい\n安心できる私的空間",
    "仕事・学校など\n義務がある公的空間",
    "どちらでもない\n「居場所」がある空間\n\n← ここを作る",
]
col_highlight = [False, False, True]
for i in range(3):
    bg = C_ACCENT if col_highlight[i] else C_WHITE
    fc = C_WHITE if col_highlight[i] else C_DARK
    card(sl, col_x[i], Inches(2.85), Inches(3.9), Inches(2.6), bg)
    add_textbox(sl, col_titles[i], col_x[i]+Inches(0.15), Inches(2.95),
                Inches(3.6), Inches(0.75), font_size=18, bold=True, color=fc)
    add_textbox(sl, col_texts[i], col_x[i]+Inches(0.15), Inches(3.75),
                Inches(3.6), Inches(1.5), font_size=14, color=fc)

add_textbox(sl,
            "理念：まごころ（見返りを求めない善意）を基盤としたコミュニティ",
            Inches(0.5), Inches(5.6), Inches(12), Inches(0.5),
            font_size=14, italic=True, color=C_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# スライド 4: なぜ必要か（課題）
# ============================================================
sl = prs.slides.add_slide(blank_layout)
bg_cream(sl)
header_bar(sl, "なぜ今、これが必要なのか", "メルボルン在住日本人が抱えるリアルな課題")
accent_line(sl)

problems = [
    ("😰 仕事が見つからない",
     "英語力・経験不足でローカル企業への応募が難しい\n求人サイトに載っていない「日本語OKな職場」の情報がない"),
    ("🏠 住居探しが不安",
     "Flatmatesなど英語プラットフォームのハードルが高い\nどのエリアが安全か・相場がわからない"),
    ("🤝 頼れる人がいない",
     "渡航直後は繋がりがゼロ\nSNSの情報は古い・不確か・悪用される危険もある"),
]

for i, (title, desc) in enumerate(problems):
    y = Inches(1.5) + i * Inches(1.7)
    card(sl, Inches(0.4), y, Inches(12.5), Inches(1.55), C_WHITE)
    add_textbox(sl, title, Inches(0.65), y+Inches(0.12),
                Inches(4), Inches(0.6), font_size=20, bold=True, color=C_ACCENT)
    add_textbox(sl, desc, Inches(0.65), y+Inches(0.62),
                Inches(11.8), Inches(0.85), font_size=15, color=C_DARK)

add_textbox(sl,
            "→ 「信頼できる先輩の経験」が最大の解決策になる",
            Inches(0.5), Inches(6.65), Inches(12), Inches(0.55),
            font_size=17, bold=True, color=C_ACCENT)

# ============================================================
# スライド 5: サービスの仕組み
# ============================================================
sl = prs.slides.add_slide(blank_layout)
bg_cream(sl)
header_bar(sl, "サービスの仕組み", "5ステップのユーザーライフサイクル")
accent_line(sl)

stages = [
    ("STAGE 0", "LINE参加", "属性把握\nアンケート\n（1分）", "-"),
    ("STAGE 1", "パッケージ\n購入", "ヒアリング\nマッチング\n開始", "A$150〜200"),
    ("STAGE 2", "到着後\n1〜2週間", "到着後\nアンケート", "A$10 CB"),
    ("STAGE 3", "滞在中\n3ヶ月ごと", "近況報告\nレビュー", "A$5/回 CB"),
    ("STAGE 4", "帰国前\n1〜2ヶ月", "引き継ぎ\nノート作成", "A$30 CB"),
]

sw = Inches(2.3)
for i, (stage, timing, action, price) in enumerate(stages):
    x = Inches(0.25) + i * Inches(2.55)
    y_top = Inches(1.5)
    # ステージ番号
    add_rect(sl, x, y_top, sw, Inches(0.45), C_ACCENT)
    add_textbox(sl, stage, x, y_top+Pt(2), sw, Inches(0.4),
                font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # カード本体
    card(sl, x, y_top+Inches(0.45), sw, Inches(4.2), C_WHITE)
    add_textbox(sl, timing, x+Inches(0.1), y_top+Inches(0.55),
                sw-Inches(0.2), Inches(0.7),
                font_size=14, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_textbox(sl, action, x+Inches(0.1), y_top+Inches(1.35),
                sw-Inches(0.2), Inches(1.5),
                font_size=13, color=C_GRAY, align=PP_ALIGN.CENTER)
    # 金額
    add_rect(sl, x, y_top+Inches(3.0), sw, Inches(0.6),
             C_SUBTLE if price == "-" else C_ACCENT)
    fc = C_GRAY if price == "-" else C_WHITE
    add_textbox(sl, price, x, y_top+Inches(3.05), sw, Inches(0.5),
                font_size=13, bold=True, color=fc, align=PP_ALIGN.CENTER)
    # 矢印
    if i < 4:
        add_textbox(sl, "→", x+sw, y_top+Inches(1.8), Inches(0.25), Inches(0.5),
                    font_size=20, bold=True, color=C_ACCENT)

add_textbox(sl,
            "12ヶ月フル参加の実質コスト：A$85〜135　（後任紹介成功で最大+A$50追加）",
            Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55),
            font_size=14, color=C_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# スライド 6: 最大の武器 ー 引き継ぎノート
# ============================================================
sl = prs.slides.add_slide(blank_layout)
bg_cream(sl)
header_bar(sl, "最大の競合優位性：引き継ぎノート", "どの求人サイト・不動産サイトにも絶対に載らない「生の声」")
accent_line(sl)

add_textbox(sl,
            "蓄積されるほど価値が高まる　→　後発の模倣が困難な参入障壁",
            Inches(0.5), Inches(1.45), Inches(12), Inches(0.5),
            font_size=16, italic=True, color=C_ACCENT)

# 左：職場版
card(sl, Inches(0.4), Inches(2.1), Inches(5.8), Inches(4.3), C_WHITE)
add_textbox(sl, "💼 職場版", Inches(0.6), Inches(2.2), Inches(5.4), Inches(0.55),
            font_size=20, bold=True, color=C_DARK)
items_w = [
    "• オーナーの名前・性格・接し方のコツ",
    "• 最初の1週間でやること",
    "• 時給・シフト・必要な英語レベル",
    "• 絶対にやってはいけないこと",
    "• 総合評価（★）",
    "• 引き継ぎ者からのメッセージ",
]
for j, t in enumerate(items_w):
    add_textbox(sl, t, Inches(0.6), Inches(2.85)+j*Inches(0.55),
                Inches(5.4), Inches(0.5), font_size=14, color=C_DARK)

# 右：住居版
card(sl, Inches(7.0), Inches(2.1), Inches(5.8), Inches(4.3), C_WHITE)
add_textbox(sl, "🏠 住居版", Inches(7.2), Inches(2.2), Inches(5.4), Inches(0.55),
            font_size=20, bold=True, color=C_DARK)
items_h = [
    "• 家賃・エリア・周辺環境",
    "• 大家・オーナーの人柄",
    "• ハウスメイト情報",
    "• 治安・交通アクセス",
    "• 光熱費・Wi-Fi等の実態",
    "• 引き継ぎ者からのメッセージ",
]
for j, t in enumerate(items_h):
    add_textbox(sl, t, Inches(7.2), Inches(2.85)+j*Inches(0.55),
                Inches(5.4), Inches(0.5), font_size=14, color=C_DARK)

# 中央の矢印スペース
add_textbox(sl, "⇆", Inches(6.0), Inches(3.8), Inches(1.0), Inches(0.8),
            font_size=30, color=C_ACCENT, align=PP_ALIGN.CENTER)

# ============================================================
# スライド 7: コミュニティの現在地
# ============================================================
sl = prs.slides.add_slide(blank_layout)
bg_cream(sl)
header_bar(sl, "コミュニティの現在地", "私たちはすでに強力な資産を持っている")
accent_line(sl)

assets = [
    ("👥", "LINEグループ\n（メルボルン在住日本人）", "300名\n（在住200名）"),
    ("📱", "Facebook\n「メルボルンにいる日本人」", "2.7万人"),
    ("📸", "Instagram\n@3rd.placejapan", "190人\n（開始直後）"),
    ("🏠", "物理拠点\n「皆の館」", "創設者マサ氏の\nアパート1室"),
    ("🤝", "パートナー店舗", "3店舗\n（直接の関係）"),
    ("💰", "確保済み投資家", "1名"),
]

col = 3
for i, (icon, label, val) in enumerate(assets):
    row, col_i = divmod(i, 3)
    x = Inches(0.4) + col_i * Inches(4.25)
    y = Inches(1.55) + row * Inches(2.4)
    card(sl, x, y, Inches(3.9), Inches(2.15), C_WHITE)
    add_textbox(sl, icon, x+Inches(0.15), y+Inches(0.12),
                Inches(0.8), Inches(0.6), font_size=28)
    add_textbox(sl, label, x+Inches(1.0), y+Inches(0.15),
                Inches(2.7), Inches(0.75), font_size=13, color=C_GRAY)
    add_textbox(sl, val, x+Inches(0.15), y+Inches(1.2),
                Inches(3.6), Inches(0.75), font_size=20, bold=True, color=C_ACCENT)

add_textbox(sl,
            "東京では月1回イベント開催（参加費500円・黒字運営）　運営参加希望者 15名",
            Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.55),
            font_size=14, color=C_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# スライド 8: ビジネスモデル
# ============================================================
sl = prs.slides.add_slide(blank_layout)
bg_cream(sl)
header_bar(sl, "ビジネスモデル", "BtoCに特化した4本柱の収益構造")
accent_line(sl)

add_textbox(sl,
            "なぜBtoC？　→　メルボルンは求職者が過多。店舗はレジュメが自然に集まるため採用費への動機がない。\n痛みを抱えているのは求職者側。だからCからお金をもらう設計。",
            Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.75),
            font_size=14, color=C_GRAY)

pillars = [
    ("🎯 メルボルン着パッケージ\n（メイン収益）",
     "A$150〜200/人",
     "・渡航前Zoomブリーフィング\n・仕事＆住居マッチング\n・レジュメ添削\n・公認推薦状\n・カレー会招待\n・到着後1ヶ月フォロー\n・キャッシュバック最大A$100"),
    ("📢 コミュニティ告知枠\n（BtoB）",
     "A$50〜200/回",
     "LINE300人・FB2.7万人への告知\n\n対象：SIM・海外保険・\n語学学校・転職エージェント等"),
    ("🔑 長期会員制",
     "A$20/月",
     "継続サポート・\nコミュニティメンバー資格\n引き継ぎノート閲覧"),
    ("🏠 住居バトンパス\n（将来展開）",
     "手数料制",
     "住居の引き継ぎ仲介\nルームメイトマッチング"),
]

for i, (title, price, detail) in enumerate(pillars):
    x = Inches(0.3) + i * Inches(3.2)
    y = Inches(2.35)
    # タイトル部
    add_rect(sl, x, y, Inches(3.0), Inches(0.75),
             C_ACCENT if i == 0 else C_DARK)
    add_textbox(sl, title, x+Inches(0.1), y+Inches(0.05),
                Inches(2.8), Inches(0.68), font_size=12, bold=True, color=C_WHITE)
    # 価格
    add_rect(sl, x, y+Inches(0.75), Inches(3.0), Inches(0.55), C_SUBTLE)
    add_textbox(sl, price, x, y+Inches(0.78), Inches(3.0), Inches(0.5),
                font_size=16, bold=True, color=C_ACCENT if i == 0 else C_DARK,
                align=PP_ALIGN.CENTER)
    # 詳細
    card(sl, x, y+Inches(1.3), Inches(3.0), Inches(3.5), C_WHITE)
    add_textbox(sl, detail, x+Inches(0.12), y+Inches(1.4),
                Inches(2.78), Inches(3.3), font_size=12, color=C_DARK)

# ============================================================
# スライド 9: フェーズ計画
# ============================================================
sl = prs.slides.add_slide(blank_layout)
bg_cream(sl)
header_bar(sl, "フェーズ計画", "段階的に信頼を積み上げてスケールへ")
accent_line(sl)

phases = [
    ("Phase 1\n〜2026年6月",
     "信頼の構築",
     "• ヒアリング 20人\n• 店舗交渉 10件\n• 初回マッチング（無料）\n• アプリβ版リリース",
     "今ここ",
     True),
    ("Phase 2\n2026年7〜9月",
     "マネタイズ開始",
     "• 有料パッケージ販売\n• BtoB告知枠営業\n• 引き継ぎノート拡充\n• SNS本格運用",
     "",
     False),
    ("Phase 3\n2026年10月〜",
     "スケール",
     "• 登録店舗 30店\n• 求職者 200人\n• マッチング 100件\n• 月収 A$2,000〜3,000",
     "",
     False),
]

kpi_rows = [
    ("", "3ヶ月目標", "12ヶ月目標"),
    ("求職者", "30人", "200人"),
    ("店舗", "5店", "30店"),
    ("マッチング", "10件", "100件"),
    ("月収", "A$200", "A$2,000〜3,000"),
]

for i, (period, theme, actions, badge, current) in enumerate(phases):
    x = Inches(0.3) + i * Inches(4.1)
    y = Inches(1.5)
    bg = C_ACCENT if current else (C_DARK if i == 2 else C_WHITE)
    fc = C_WHITE if current or i == 2 else C_DARK
    fc2 = C_WHITE if current or i == 2 else C_GRAY
    card(sl, x, y, Inches(3.8), Inches(4.85), bg)
    if current:
        add_textbox(sl, "▶ 現在", x+Inches(2.3), y+Inches(0.12),
                    Inches(1.3), Inches(0.35), font_size=11, bold=True, color=C_WHITE)
    add_textbox(sl, period, x+Inches(0.15), y+Inches(0.1),
                Inches(3.5), Inches(0.75), font_size=16, bold=True, color=fc)
    add_textbox(sl, theme, x+Inches(0.15), y+Inches(0.85),
                Inches(3.5), Inches(0.5), font_size=18, bold=True, color=fc)
    add_textbox(sl, actions, x+Inches(0.15), y+Inches(1.45),
                Inches(3.5), Inches(3.2), font_size=13, color=fc2)

# KPIテーブル
y_t = Inches(6.45)
col_ws = [Inches(3.5), Inches(2.5), Inches(2.5)]
col_xs = [Inches(2.8), Inches(6.5), Inches(9.2)]
for row_i, row in enumerate(kpi_rows):
    for col_i, cell in enumerate(row):
        bg = C_DARK if row_i == 0 else (C_SUBTLE if row_i % 2 == 0 else C_WHITE)
        fc = C_WHITE if row_i == 0 else C_DARK
        add_rect(sl, col_xs[col_i], y_t + row_i * Inches(0.2),
                 col_ws[col_i], Inches(0.2), bg)
        # テキストは省略（スペース不足のため次スライドへ）

# ============================================================
# スライド 10: ビジョン
# ============================================================
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, C_DARK)
add_rect(sl, 0, Inches(5.5), SLIDE_W, Inches(2.0), C_ACCENT)

add_textbox(sl, "私たちが目指す世界",
            Inches(1), Inches(0.8), Inches(11), Inches(0.8),
            font_size=20, color=C_ACCENT, align=PP_ALIGN.CENTER)

add_textbox(sl,
            "「メルボルンに行くなら\n3rd Place Hub に相談する」",
            Inches(0.5), Inches(1.7), Inches(12.3), Inches(2.0),
            font_size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_textbox(sl,
            "が当たり前になるコミュニティ",
            Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.8),
            font_size=28, color=C_LIGHT, align=PP_ALIGN.CENTER)

vision_items = [
    "🔄  先輩から後輩へ。経験が循環するエコシステム",
    "🌏  メルボルン発・海外在住日本人の新しいセーフティネット",
    "💚  まごころが報われるコミュニティ経済",
]
for i, txt in enumerate(vision_items):
    add_textbox(sl, txt,
                Inches(1.5), Inches(5.6) + i * Inches(0.55),
                Inches(10.3), Inches(0.5),
                font_size=15, bold=(i == 0), color=C_DARK)

# ============================================================
# スライド 11: 一緒にやりましょう
# ============================================================
sl = prs.slides.add_slide(blank_layout)
bg_cream(sl)
header_bar(sl, "一緒に運営しませんか", "3rd Place Hub は仲間を探しています")
accent_line(sl)

add_textbox(sl,
            "私たちが求めているのは、メルボルンという場所が好きで、\n誰かの役に立つことに喜びを感じる人です。",
            Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.8),
            font_size=17, color=C_DARK, align=PP_ALIGN.CENTER)

roles = [
    ("🤝 マッチング担当", "求職者のヒアリング\n店舗との関係構築\n引き継ぎノートの整備"),
    ("📱 SNS・発信担当", "Instagram/LINE運用\n体験談コンテンツ制作\nコミュニティイベント"),
    ("💼 営業・渉外担当", "店舗開拓・交渉\nBtoB広告枠の提案\n投資家向け資料"),
    ("💻 テック担当", "アプリ開発・改善\nDB設計・自動化\nKPI管理"),
]

for i, (role, detail) in enumerate(roles):
    x = Inches(0.3) + (i % 2) * Inches(6.4)
    y = Inches(2.5) + (i // 2) * Inches(1.9)
    card(sl, x, y, Inches(6.1), Inches(1.7), C_WHITE)
    add_textbox(sl, role, x+Inches(0.15), y+Inches(0.1),
                Inches(5.8), Inches(0.55), font_size=18, bold=True, color=C_ACCENT)
    add_textbox(sl, detail, x+Inches(0.15), y+Inches(0.65),
                Inches(5.8), Inches(0.95), font_size=13, color=C_DARK)

add_textbox(sl,
            "まずは一緒にカレー会から！　Instagram: @3rd.placejapan",
            Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.55),
            font_size=15, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# ============================================================
# 保存
# ============================================================
out_path = "/Users/fukushidaichi/Desktop/3rd place/3rd_place_website/3rdplace_hub_presentation.pptx"
prs.save(out_path)
print(f"✅ 保存完了: {out_path}")
